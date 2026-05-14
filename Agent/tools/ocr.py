"""
OCR 与文件处理工具
- GPT-4o 多模态 OCR（图片文字识别）
- PDF/Word/PPT 文本提取
- LLM 自动分类
"""
import base64
import uuid
from pathlib import Path
from typing import Optional, Dict, Tuple
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage
from config import OPENAI_API_KEY, OPENAI_BASE_URL, UPLOAD_DIR


def _get_openai_chat(**kwargs) -> ChatOpenAI:
    """创建 ChatOpenAI 实例，自动注入 base_url"""
    if OPENAI_BASE_URL:
        kwargs["base_url"] = OPENAI_BASE_URL
    return ChatOpenAI(api_key=OPENAI_API_KEY, **kwargs)


def _encode_image(image_path: str) -> str:
    """图片文件 → base64 data URL"""
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def ocr_image(image_path: str) -> str:
    """
    使用 GPT-4o 进行图片 OCR 识别
    支持手写体、印刷体中文和数学公式
    """
    if not OPENAI_API_KEY:
        raise RuntimeError("OPENAI_API_KEY not configured")

    llm = _get_openai_chat(model="Doubao-Seed-2.0-pro", max_tokens=4096)
    b64 = _encode_image(image_path)
    ext = Path(image_path).suffix.lower().replace(".", "")

    msg = HumanMessage(content=[
        {
            "type": "text",
            "text": (
                "请精确识别并转录这张图片中的所有文字内容。"
                "要求：\n"
                "1. 保留原文的段落结构和标题层级\n"
                "2. 数学公式用 LaTeX 格式表示\n"
                "3. 表格保留行列结构，用 Markdown 表格格式\n"
                "4. 如果文字模糊，标注 [不确定: 猜测内容]\n"
                "5. 不要添加任何多余的解释或开场白，直接输出识别结果"
            ),
        },
        {
            "type": "image_url",
            "image_url": {"url": f"data:image/{ext};base64,{b64}"},
        },
    ])

    response = llm.invoke([msg])
    return response.content


def extract_pdf_text(file_path: str) -> str:
    """提取 PDF 文本（文字型 PDF）"""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    texts = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text()
        if t:
            texts.append(f"[第 {i+1} 页]\n{t}")
    return "\n\n".join(texts)


def extract_docx_text(file_path: str) -> str:
    """提取 Word 文档文本"""
    from docx import Document
    doc = Document(file_path)
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    return "\n\n".join(texts)


def extract_pptx_text(file_path: str) -> str:
    """提取 PPT 文本"""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
        if slide_texts:
            texts.append(f"[幻灯片 {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def extract_text_from_file(file_path: str) -> Tuple[str, str]:
    """
    根据文件类型提取文本
    返回: (file_type, extracted_text)
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"):
        return "image", ocr_image(file_path)
    elif ext == ".pdf":
        text = extract_pdf_text(file_path)
        if len(text.strip()) < 50:
            # 可能是扫描型 PDF，暂时返回提取结果
            # 扫描型 PDF 需要先转图片再 OCR，MVP 版本简化处理
            return "pdf_scanned", text
        return "pdf", text
    elif ext in (".docx", ".doc"):
        return "word", extract_docx_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return "ppt", extract_pptx_text(file_path)
    elif ext in (".txt", ".md", ".markdown"):
        with open(file_path, "r", encoding="utf-8") as f:
            return "text", f.read()
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def classify_material(text_preview: str) -> Dict[str, str]:
    """
    使用 LLM 对资料进行自动分类
    取文本前 500 字，返回 {科目, 章节, 知识点}
    """
    if not OPENAI_API_KEY:
        return {"subject": "未分类", "chapter": "", "knowledge_point": ""}

    llm = _get_openai_chat(model="Doubao-Seed-2.0-pro", max_tokens=256)
    truncated = text_preview[:500]

    prompt = f"""请分析以下学习资料片段，输出分类结果。严格按照 JSON 格式返回，不要有其他内容。

资料内容：
{truncated}

返回格式：
{{"subject": "科目名称（如 数学/语文/英语/物理/化学/历史/政治/生物/地理/计算机/其他）",
 "chapter": "所属章节（如 函数与导数/力学/有机化学）",
 "knowledge_point": "具体知识点（如 二次函数单调性/牛顿第二定律/定语从句）"}}

注意：如果无法判断某个字段，填写"未分类"。"""

    response = llm.invoke([HumanMessage(content=prompt)])

    import json
    try:
        # 尝试从响应中提取 JSON
        text = response.content.strip()
        if "```json" in text:
            text = text.split("```json")[1].split("```")[0].strip()
        elif "```" in text:
            text = text.split("```")[1].split("```")[0].strip()
        return json.loads(text)
    except (json.JSONDecodeError, IndexError):
        return {"subject": "未分类", "chapter": "", "knowledge_point": ""}


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """
    文本分块
    按自然段落分割，尽量保持段落完整
    """
    paragraphs = text.split("\n\n")
    chunks = []
    current = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) < chunk_size:
            current = (current + "\n\n" + para).strip()
        else:
            if current:
                chunks.append(current)
            # 如果单段超过 chunk_size，按句子进一步切分
            if len(para) > chunk_size:
                sentences = para.replace("。", "。\n").split("\n")
                sub = ""
                for s in sentences:
                    s = s.strip()
                    if not s:
                        continue
                    if len(sub) + len(s) < chunk_size:
                        sub = (sub + s).strip()
                    else:
                        if sub:
                            chunks.append(sub)
                        sub = s
                if sub:
                    current = sub
            else:
                current = para

    if current:
        chunks.append(current)

    return chunks
